from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Colours ───────────────────────────────────────────────
RED    = RGBColor(0xD6, 0x28, 0x28)
NAVY   = RGBColor(0x1a, 0x1a, 0x2e)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LGREY  = RGBColor(0xF4, 0xF5, 0xF7)
MGREY  = RGBColor(0x55, 0x55, 0x55)
BGREY  = RGBColor(0xAA, 0xAA, 0xAA)
BLUE   = RGBColor(0x4a, 0x90, 0xD9)
GREEN  = RGBColor(0x5c, 0xb8, 0x5c)
ORANGE = RGBColor(0xf0, 0xad, 0x4e)
DGREY  = RGBColor(0x14, 0x14, 0x25)
TEAL   = RGBColor(0x1a, 0x8f, 0x8f)

W = Inches(13.33)
H = Inches(7.50)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
blank = prs.slide_layouts[6]


# ── Helpers ───────────────────────────────────────────────

def add_slide():
    return prs.slides.add_slide(blank)

def rect(slide, left, top, width, height, fill):
    s = slide.shapes.add_shape(1, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    return s

def txt(slide, text, left, top, width, height,
        size=16, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    box.word_wrap = True
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size  = Pt(size)
    r.font.bold  = bold
    r.font.color.rgb = color
    r.font.italic = italic

def tag(slide, label, left, top):
    rect(slide, left, top, Inches(1.9), Inches(0.35), RED)
    txt(slide, label, left + Inches(0.1), top + Inches(0.05),
        Inches(1.7), Inches(0.28), size=10, bold=True,
        color=WHITE, align=PP_ALIGN.CENTER)

def slide_number(slide, n, total=4):
    txt(slide, f"{n} / {total}",
        W - Inches(1.1), Inches(0.12), Inches(0.95), Inches(0.3),
        size=10, color=BGREY, align=PP_ALIGN.RIGHT)

def bullet_list(slide, items, left, top, width,
                text_color=NAVY, dot_color=RED, size=14):
    y = top
    for item in items:
        d = slide.shapes.add_shape(1, left, y + Inches(0.14),
                                   Inches(0.1), Inches(0.1))
        d.fill.solid(); d.fill.fore_color.rgb = dot_color
        d.line.fill.background()
        txt(slide, item, left + Inches(0.22), y,
            width - Inches(0.25), Inches(0.48),
            size=size, color=text_color)
        y += Inches(0.52)

def red_line(slide, top):
    rect(slide, Inches(0.45), top, Inches(0.5), Inches(0.06), RED)


# ═══════════════════════════════════════════════════════════
# SLIDE 1 — PROBLEM
# ═══════════════════════════════════════════════════════════
s1 = add_slide()

LW = Inches(5.6)
rect(s1, 0, 0, LW, H, WHITE)
rect(s1, LW, 0, W - LW, H, LGREY)
rect(s1, 0, 0, LW, Inches(0.08), RED)

slide_number(s1, 1)
tag(s1, "THE PROBLEM", Inches(0.45), Inches(0.7))

txt(s1, "Refugees in Switzerland",
    Inches(0.45), Inches(1.3), Inches(4.9), Inches(0.6),
    size=30, bold=True, color=NAVY)
txt(s1, "don't know their rights.",
    Inches(0.45), Inches(1.85), Inches(4.9), Inches(0.6),
    size=30, bold=True, color=RED)

red_line(s1, Inches(2.55))

txt(s1, "Switzerland's asylum system is one of the most complex in Europe — "
        "5 permit types, canton-specific rules, 30-day appeal deadlines, "
        "and official information only in German, French, or Italian.",
    Inches(0.45), Inches(2.75), Inches(4.9), Inches(1.1),
    size=13, color=MGREY)

stats = [
    ("100K+", "Refugees & asylum\nseekers in Switzerland"),
    ("5",     "Permit types, each with\ndifferent rights & rules"),
    ("0",     "Free multilingual tools\nexplaining it clearly"),
]
sx = Inches(0.45)
for num, lbl in stats:
    rect(s1, sx, Inches(4.1), Inches(1.55), Inches(1.05),
         RGBColor(0xFF, 0xF5, 0xF5))
    rect(s1, sx, Inches(4.1), Inches(0.04), Inches(1.05), RED)
    txt(s1, num, sx + Inches(0.1), Inches(4.16),
        Inches(1.4), Inches(0.45), size=24, bold=True, color=RED)
    txt(s1, lbl, sx + Inches(0.1), Inches(4.58),
        Inches(1.4), Inches(0.5), size=9, color=MGREY)
    sx += Inches(1.65)

txt(s1, '"Missed deadlines. Lost rights. Fear and isolation.\nNot from bad law — from inaccessible information."',
    Inches(0.45), Inches(5.4), Inches(4.9), Inches(0.9),
    size=12, italic=True, color=MGREY)

# Right panel — Swiss cross illustration
cx = LW + (W - LW) / 2
cy = H / 2
rect(s1, cx - Inches(0.55), cy - Inches(2.0), Inches(1.1), Inches(4.0),
     RGBColor(0xF0, 0xD0, 0xD0))
rect(s1, cx - Inches(2.0), cy - Inches(0.55), Inches(4.0), Inches(1.1),
     RGBColor(0xF0, 0xD0, 0xD0))
rect(s1, cx - Inches(0.5), cy - Inches(1.15), Inches(1.0), Inches(1.0),
     RGBColor(0x3d, 0x3d, 0x5c))
rect(s1, cx - Inches(0.45), cy - Inches(0.2), Inches(0.9), Inches(1.2),
     RGBColor(0x3d, 0x3d, 0x5c))
txt(s1, "?", cx - Inches(2.0), cy - Inches(0.4),
    Inches(0.8), Inches(0.9), size=60, bold=True,
    color=RGBColor(0xD6, 0x28, 0x28), align=PP_ALIGN.CENTER)
txt(s1, "?", cx + Inches(1.3), cy - Inches(0.4),
    Inches(0.8), Inches(0.9), size=60, bold=True,
    color=RGBColor(0xD6, 0x28, 0x28), align=PP_ALIGN.CENTER)
txt(s1, "?", cx - Inches(0.35), cy - Inches(2.0),
    Inches(0.8), Inches(0.9), size=44, bold=True,
    color=RGBColor(0xD6, 0x28, 0x28), align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════
# SLIDE 2 — SOLUTION
# ═══════════════════════════════════════════════════════════
s2 = add_slide()
rect(s2, 0, 0, LW, H, NAVY)
rect(s2, LW, 0, W - LW, H, DGREY)
rect(s2, 0, 0, LW, Inches(0.08), RED)

slide_number(s2, 2)
tag(s2, "THE SOLUTION", Inches(0.45), Inches(0.7))

txt(s2, "AmanCH",
    Inches(0.45), Inches(1.3), Inches(4.9), Inches(0.65),
    size=38, bold=True, color=WHITE)
txt(s2, "امن  ·  Peace in Arabic. CH  ·  Switzerland.",
    Inches(0.45), Inches(1.95), Inches(4.9), Inches(0.38),
    size=12, italic=True, color=BGREY)

txt(s2, "A free AI assistant available 24/7 — answering questions about\n"
        "asylum, work, education, healthcare, and integration\n"
        "in the refugee's own language.",
    Inches(0.45), Inches(2.5), Inches(4.9), Inches(1.0),
    size=13, color=BGREY)

red_line(s2, Inches(3.6))

bullet_list(s2, [
    "15+ languages — Arabic, Tigrinya, Dari, Somali, Ukrainian, French…",
    "Covers work rights, school, health, permits, family reunification",
    "Grounded in 70+ verified Swiss official sources (SEM, OSAR, cantons)",
    "Always on — no appointment, no waiting, no account needed",
    "Proactively warns about the critical 30-day appeal deadline",
    "Free for every refugee in Switzerland",
], Inches(0.45), Inches(3.8), Inches(4.9),
   text_color=WHITE, dot_color=RED, size=13)

# Right — phone mockup
px = LW + Inches(1.3)
py = Inches(0.45)
pw = Inches(3.7)
ph = Inches(6.4)

rect(s2, px, py, pw, ph, RGBColor(0x2a, 0x2a, 0x3e))
rect(s2, px + Inches(0.14), py + Inches(0.22),
     pw - Inches(0.28), ph - Inches(0.35), DGREY)

rect(s2, px + Inches(0.14), py + Inches(0.65),
     pw - Inches(0.28), Inches(0.7), RED)
txt(s2, "AmanCH",
    px + Inches(0.22), py + Inches(0.7),
    pw - Inches(0.4), Inches(0.32),
    size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s2, "Online · Replies in seconds",
    px + Inches(0.22), py + Inches(1.02),
    pw - Inches(0.4), Inches(0.25),
    size=8, color=BGREY, align=PP_ALIGN.CENTER)

chats = [
    ("user", "What is Permit F?"),
    ("bot",  "Permit F = provisionally admitted.\nYou may work with cantonal authorisation."),
    ("user", "هل يمكنني العمل؟"),
    ("bot",  "نعم، يمكنك العمل بعد الحصول على إذن."),
    ("user", "Puis-je inscrire mon enfant à l'école?"),
    ("bot",  "Oui, tous les enfants ont droit à l'école\nen Suisse, quel que soit le permis."),
]
cy_c = py + Inches(1.5)
for role, msg in chats:
    if role == "user":
        bx2 = px + pw - Inches(2.7)
        bc2, tc2 = RED, WHITE
    else:
        bx2 = px + Inches(0.22)
        bc2, tc2 = RGBColor(0x2a, 0x2a, 0x3e), BGREY
    bh2 = Inches(0.52) if "\n" not in msg else Inches(0.72)
    rect(s2, bx2, cy_c, Inches(2.35), bh2, bc2)
    txt(s2, msg, bx2 + Inches(0.08), cy_c + Inches(0.06),
        Inches(2.2), bh2, size=8, color=tc2)
    cy_c += bh2 + Inches(0.1)


# ═══════════════════════════════════════════════════════════
# SLIDE 3 — DEMO (Live Product)
# ═══════════════════════════════════════════════════════════
s3 = add_slide()
rect(s3, 0, 0, W, H, WHITE)
rect(s3, 0, 0, W, Inches(0.08), RED)
rect(s3, 0, Inches(1.35), W, Inches(0.03), LGREY)

slide_number(s3, 3)
tag(s3, "LIVE DEMO", Inches(0.45), Inches(0.38))
txt(s3, "AmanCH — Live, Deployed, Used Today",
    Inches(2.6), Inches(0.32), Inches(9.5), Inches(0.7),
    size=26, bold=True, color=NAVY, align=PP_ALIGN.LEFT)

# Left column — browser mockup
bm_x = Inches(0.45)
bm_y = Inches(1.55)
bm_w = Inches(7.2)
bm_h = Inches(5.55)

rect(s3, bm_x, bm_y, bm_w, bm_h, RGBColor(0xE8, 0xE8, 0xE8))
rect(s3, bm_x, bm_y, bm_w, Inches(0.45), RGBColor(0xD0, 0xD0, 0xD0))

# Browser dots
for di, dc in enumerate([RGBColor(0xFF, 0x5F, 0x56),
                          RGBColor(0xFF, 0xBD, 0x2E),
                          RGBColor(0x27, 0xC9, 0x3F)]):
    dot = s3.shapes.add_shape(1,
        bm_x + Inches(0.15) + di * Inches(0.25),
        bm_y + Inches(0.14), Inches(0.15), Inches(0.15))
    dot.fill.solid(); dot.fill.fore_color.rgb = dc
    dot.line.fill.background()

# URL bar
rect(s3, bm_x + Inches(0.7), bm_y + Inches(0.08),
     bm_w - Inches(0.85), Inches(0.28), WHITE)
txt(s3, "refugee-assistant-switzerland-gndat6oqniq737kahzd3p9.streamlit.app",
    bm_x + Inches(0.75), bm_y + Inches(0.1),
    bm_w - Inches(0.9), Inches(0.22),
    size=8, color=MGREY)

# App inside browser
app_y = bm_y + Inches(0.5)
app_h = bm_h - Inches(0.55)
rect(s3, bm_x, app_y, bm_w, app_h, NAVY)

# App header bar
rect(s3, bm_x, app_y, bm_w, Inches(0.6), RED)
txt(s3, "AmanCH — Refugee Assistant Switzerland",
    bm_x + Inches(0.2), app_y + Inches(0.12),
    bm_w - Inches(0.4), Inches(0.35),
    size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Permit selector strip
rect(s3, bm_x, app_y + Inches(0.6), bm_w, Inches(0.38),
     RGBColor(0x22, 0x22, 0x42))
permits = ["N", "F", "B", "C", "S", "?"]
px2 = bm_x + Inches(0.3)
for p in permits:
    bg = RED if p == "F" else RGBColor(0x33, 0x33, 0x55)
    rect(s3, px2, app_y + Inches(0.65), Inches(0.55), Inches(0.25), bg)
    txt(s3, p, px2, app_y + Inches(0.65),
        Inches(0.55), Inches(0.25),
        size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    px2 += Inches(0.65)

# Chat messages inside app
demo_chats = [
    ("user", "Permit F",   "Can I work in Switzerland?"),
    ("bot",  "AmanCH",     "Yes — with Permit F you can work after receiving\ncantonal authorisation from your migration office.\nSource: sem.admin.ch"),
    ("user", "Permit F",   "كم يستغرق طلب تصريح العمل؟"),
    ("bot",  "AmanCH",     "عادةً ما يستغرق القرار من أسبوع إلى ثلاثة أسابيع\nحسب الكانتون. تواصل مع مكتب الهجرة في كانتونك."),
]
chat_y = app_y + Inches(1.1)
for role, sender, msg in demo_chats:
    if role == "user":
        cbx = bm_x + bm_w - Inches(3.8)
        cbc = RGBColor(0x8B, 0x1A, 0x1A)
        ctc = WHITE
    else:
        cbx = bm_x + Inches(0.2)
        cbc = RGBColor(0x2a, 0x2a, 0x4e)
        ctc = BGREY
    lines = msg.count("\n") + 1
    cbh = Inches(0.38 + 0.28 * lines)
    rect(s3, cbx, chat_y, Inches(3.6), cbh, cbc)
    txt(s3, msg, cbx + Inches(0.1), chat_y + Inches(0.06),
        Inches(3.4), cbh, size=9, color=ctc)
    chat_y += cbh + Inches(0.12)

# Input bar at bottom of app
rect(s3, bm_x, app_y + app_h - Inches(0.5), bm_w, Inches(0.5),
     RGBColor(0x22, 0x22, 0x42))
rect(s3, bm_x + Inches(0.2), app_y + app_h - Inches(0.42),
     bm_w - Inches(1.2), Inches(0.32), RGBColor(0x33, 0x33, 0x55))
txt(s3, "Type your question in any language…",
    bm_x + Inches(0.3), app_y + app_h - Inches(0.4),
    bm_w - Inches(1.4), Inches(0.28),
    size=9, italic=True, color=BGREY)
rect(s3, bm_x + bm_w - Inches(0.85), app_y + app_h - Inches(0.44),
     Inches(0.6), Inches(0.34), RED)
txt(s3, "→", bm_x + bm_w - Inches(0.85), app_y + app_h - Inches(0.44),
    Inches(0.6), Inches(0.34),
    size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Right column — key stats
col_x = bm_x + bm_w + Inches(0.35)
col_w = W - col_x - Inches(0.3)

stats3 = [
    ("70+",  "Verified Swiss\nofficial sources"),
    ("15+",  "Languages\nsupported"),
    ("5",    "Permit types\ncovered"),
    ("24/7", "Available, free,\nno account needed"),
]
sy = Inches(1.55)
for val, lbl in stats3:
    rect(s3, col_x, sy, col_w, Inches(1.35), RGBColor(0xFF, 0xF5, 0xF5))
    rect(s3, col_x, sy, col_w, Inches(0.05), RED)
    txt(s3, val, col_x, sy + Inches(0.1),
        col_w, Inches(0.6), size=30, bold=True,
        color=RED, align=PP_ALIGN.CENTER)
    txt(s3, lbl, col_x, sy + Inches(0.7),
        col_w, Inches(0.55), size=10,
        color=MGREY, align=PP_ALIGN.CENTER)
    sy += Inches(1.45)

# Live URL badge
rect(s3, col_x, sy + Inches(0.1), col_w, Inches(0.55), NAVY)
txt(s3, "LIVE NOW",
    col_x, sy + Inches(0.15), col_w, Inches(0.25),
    size=9, bold=True, color=RED, align=PP_ALIGN.CENTER)
txt(s3, "streamlit.app",
    col_x, sy + Inches(0.35), col_w, Inches(0.25),
    size=8, color=BGREY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════
# SLIDE 4 — FUTURE ROADMAP
# ═══════════════════════════════════════════════════════════
s4 = add_slide()
rect(s4, 0, 0, W, H, WHITE)
rect(s4, 0, 0, W, Inches(0.08), RED)
rect(s4, 0, Inches(1.35), W, Inches(0.03), LGREY)

slide_number(s4, 4)
tag(s4, "FUTURE ROADMAP", Inches(0.45), Inches(0.38))
txt(s4, "Where AmanCH Goes Next",
    Inches(2.65), Inches(0.32), Inches(9.0), Inches(0.7),
    size=26, bold=True, color=NAVY, align=PP_ALIGN.LEFT)

# Timeline spine
spine_y = Inches(3.1)
rect(s4, Inches(0.6), spine_y, W - Inches(1.2), Inches(0.06), LGREY)

phases = [
    {
        "label": "PHASE 1",
        "period": "Now — MVP Complete",
        "color": GREEN,
        "icon": "✓",
        "title": "Foundation Built",
        "items": [
            "AI chatbot live on Streamlit Cloud",
            "70+ Swiss official sources loaded",
            "15+ languages, voice input & TTS",
            "All 5 permit types supported",
            "30-day appeal deadline alerts",
            "Free, no account needed",
        ],
        "done": True,
    },
    {
        "label": "PHASE 2",
        "period": "3–6 Months",
        "color": BLUE,
        "icon": "▶",
        "title": "Scale & Deepen",
        "items": [
            "React frontend — mobile-first UI",
            "Admin dashboard for NGOs",
            "All 26 cantons with local data",
            "Document upload & guidance",
            "Offline mode for low-connectivity",
            "Lawyer referral integration",
        ],
        "done": False,
    },
    {
        "label": "PHASE 3",
        "period": "6–12 Months",
        "color": ORANGE,
        "icon": "★",
        "title": "Impact at Scale",
        "items": [
            "Native iOS & Android app",
            "NGO & cantonal office partnerships",
            "Expand to Germany, Austria, France",
            "Real-time SEM decision tracking",
            "Anonymous usage analytics for policy",
            "Grant funding & sustainability model",
        ],
        "done": False,
    },
]

phase_w = Inches(3.8)
phase_gap = Inches(0.37)
total_pw = 3 * phase_w + 2 * phase_gap
phase_x0 = (W - total_pw) / 2

for i, ph in enumerate(phases):
    px3 = phase_x0 + i * (phase_w + phase_gap)
    col = ph["color"]

    # Dot on spine
    dot_cx = px3 + phase_w / 2 - Inches(0.18)
    dot = s4.shapes.add_shape(1, dot_cx, spine_y - Inches(0.15),
                               Inches(0.36), Inches(0.36))
    dot.fill.solid(); dot.fill.fore_color.rgb = col
    dot.line.fill.background()

    # Phase label + period above spine
    txt(s4, ph["label"],
        px3, spine_y - Inches(0.72), phase_w, Inches(0.3),
        size=10, bold=True, color=col, align=PP_ALIGN.CENTER)
    txt(s4, ph["period"],
        px3, spine_y - Inches(0.42), phase_w, Inches(0.28),
        size=9, color=MGREY, align=PP_ALIGN.CENTER)

    # Card below spine
    card_y = spine_y + Inches(0.32)
    card_h = Inches(3.65)
    rect(s4, px3, card_y, phase_w, card_h, LGREY)
    rect(s4, px3, card_y, phase_w, Inches(0.07), col)

    # Card header
    rect(s4, px3, card_y + Inches(0.07), phase_w, Inches(0.62), col)
    txt(s4, ph["icon"] + "  " + ph["title"],
        px3 + Inches(0.15), card_y + Inches(0.16),
        phase_w - Inches(0.3), Inches(0.42),
        size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Bullet items
    item_y = card_y + Inches(0.82)
    for item in ph["items"]:
        d2 = s4.shapes.add_shape(1,
            px3 + Inches(0.2), item_y + Inches(0.14),
            Inches(0.09), Inches(0.09))
        d2.fill.solid(); d2.fill.fore_color.rgb = col
        d2.line.fill.background()
        item_color = MGREY if ph["done"] else NAVY
        txt(s4, item,
            px3 + Inches(0.37), item_y,
            phase_w - Inches(0.5), Inches(0.44),
            size=11, color=item_color)
        item_y += Inches(0.46)

# Bottom vision statement
rect(s4, Inches(0.6), Inches(7.0), W - Inches(1.2), Inches(0.38),
     RGBColor(0xFF, 0xF5, 0xF5))
rect(s4, Inches(0.6), Inches(7.0), W - Inches(1.2), Inches(0.04), RED)
txt(s4, "Vision: Every refugee in Switzerland — and Europe — has a trusted, knowledgeable guide in their own language, in their pocket, for free.",
    Inches(0.8), Inches(7.04), W - Inches(1.6), Inches(0.32),
    size=11, italic=True, color=RED, align=PP_ALIGN.CENTER)


# ── Save ─────────────────────────────────────────────────
out = "AmanCH_PitchDeck.pptx"
prs.save(out)
print(f"Saved: {out}")
