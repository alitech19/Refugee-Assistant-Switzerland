from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Colours ───────────────────────────────────────────────
RED    = RGBColor(0xD6, 0x28, 0x28)
NAVY   = RGBColor(0x1a, 0x1a, 0x2e)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LGREY  = RGBColor(0xF4, 0xF5, 0xF7)
MGREY  = RGBColor(0x55, 0x55, 0x55)
BGREY  = RGBColor(0xAA, 0xAA, 0xAA)
BLUE   = RGBColor(0x4a, 0x90, 0xD9)
GREEN  = RGBColor(0x5c, 0xb8, 0x5c)
ORANGE = RGBColor(0xf0, 0xad, 0x4e)
DGREY  = RGBColor(0x14, 0x14, 0x25)

W = Inches(13.33)
H = Inches(7.50)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
blank = prs.slide_layouts[6]


# ── Helpers ───────────────────────────────────────────────

def add_slide():
    return prs.slides.add_slide(blank)

def rect(slide, left, top, width, height, fill):
    s = slide.shapes.add_shape(1, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    return s

def txt(slide, text, left, top, width, height,
        size=16, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    box.word_wrap = True
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size  = Pt(size)
    r.font.bold  = bold
    r.font.color.rgb = color
    r.font.italic = italic

def tag(slide, label, left, top):
    rect(slide, left, top, Inches(1.8), Inches(0.35), RED)
    txt(slide, label, left + Inches(0.1), top + Inches(0.05),
        Inches(1.6), Inches(0.28), size=10, bold=True,
        color=WHITE, align=PP_ALIGN.CENTER)

def slide_number(slide, n, total=4):
    txt(slide, f"{n} / {total}",
        W - Inches(1.1), Inches(0.12), Inches(0.95), Inches(0.3),
        size=10, color=BGREY, align=PP_ALIGN.RIGHT)

def bullet_list(slide, items, left, top, width,
                text_color=NAVY, dot_color=RED, size=14):
    y = top
    for item in items:
        d = slide.shapes.add_shape(1, left, y + Inches(0.14),
                                   Inches(0.1), Inches(0.1))
        d.fill.solid(); d.fill.fore_color.rgb = dot_color
        d.line.fill.background()
        txt(slide, item, left + Inches(0.22), y,
            width - Inches(0.25), Inches(0.48),
            size=size, color=text_color)
        y += Inches(0.52)

def red_line(slide, top):
    r = rect(slide, Inches(0.45), top, Inches(0.5), Inches(0.06), RED)


# ═══════════════════════════════════════════════════════════
# SLIDE 1 — PROBLEM
# ═══════════════════════════════════════════════════════════
s1 = add_slide()

# Background: white left, light grey right
LW = Inches(5.6)
rect(s1, 0, 0, LW, H, WHITE)
rect(s1, LW, 0, W - LW, H, LGREY)

# Red top accent bar
rect(s1, 0, 0, LW, Inches(0.08), RED)

slide_number(s1, 1)
tag(s1, "THE PROBLEM", Inches(0.45), Inches(0.7))

txt(s1, "Refugees in Switzerland",
    Inches(0.45), Inches(1.3), Inches(4.9), Inches(0.6),
    size=30, bold=True, color=NAVY)
txt(s1, "don't know their rights.",
    Inches(0.45), Inches(1.85), Inches(4.9), Inches(0.6),
    size=30, bold=True, color=RED)

red_line(s1, Inches(2.55))

txt(s1, "Switzerland's asylum system is one of the most complex in Europe — "
        "5 permit types, canton-specific rules, 30-day appeal deadlines, "
        "and official information only in German, French, or Italian.",
    Inches(0.45), Inches(2.75), Inches(4.9), Inches(1.1),
    size=13, color=MGREY)

# Stats
stats = [
    ("100K+", "Refugees & asylum\nseekers in Switzerland"),
    ("5",     "Permit types, each with\ndifferent rights & rules"),
    ("0",     "Free multilingual tools\nexplaining it clearly"),
]
sx = Inches(0.45)
for num, lbl in stats:
    rect(s1, sx, Inches(4.1), Inches(1.55), Inches(1.05),
         RGBColor(0xFF, 0xF5, 0xF5))
    rect(s1, sx, Inches(4.1), Inches(0.04), Inches(1.05), RED)
    txt(s1, num, sx + Inches(0.1), Inches(4.16),
        Inches(1.4), Inches(0.45), size=24, bold=True, color=RED)
    txt(s1, lbl, sx + Inches(0.1), Inches(4.58),
        Inches(1.4), Inches(0.5), size=9, color=MGREY)
    sx += Inches(1.65)

txt(s1, '"Missed deadlines. Lost rights. Fear and isolation.\nNot from bad law — from inaccessible information."',
    Inches(0.45), Inches(5.4), Inches(4.9), Inches(0.9),
    size=12, italic=True, color=MGREY)

# Right panel illustration
cx = LW + (W - LW) / 2
cy = H / 2

# Swiss cross (faded)
rect(s1, cx - Inches(0.55), cy - Inches(2.0), Inches(1.1), Inches(4.0),
     RGBColor(0xF0, 0xD0, 0xD0))
rect(s1, cx - Inches(2.0), cy - Inches(0.55), Inches(4.0), Inches(1.1),
     RGBColor(0xF0, 0xD0, 0xD0))

# Person
rect(s1, cx - Inches(0.5), cy - Inches(1.15), Inches(1.0), Inches(1.0),
     RGBColor(0x3d, 0x3d, 0x5c))
rect(s1, cx - Inches(0.45), cy - Inches(0.2), Inches(0.9), Inches(1.2),
     RGBColor(0x3d, 0x3d, 0x5c))

# Question marks
txt(s1, "?", cx - Inches(2.0), cy - Inches(0.4),
    Inches(0.8), Inches(0.9), size=60, bold=True,
    color=RGBColor(0xD6, 0x28, 0x28), align=PP_ALIGN.CENTER)
txt(s1, "?", cx + Inches(1.3), cy - Inches(0.4),
    Inches(0.8), Inches(0.9), size=60, bold=True,
    color=RGBColor(0xD6, 0x28, 0x28), align=PP_ALIGN.CENTER)
txt(s1, "?", cx - Inches(0.35), cy - Inches(2.0),
    Inches(0.8), Inches(0.9), size=44, bold=True,
    color=RGBColor(0xD6, 0x28, 0x28), align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════
# SLIDE 2 — SOLUTION
# ═══════════════════════════════════════════════════════════
s2 = add_slide()
rect(s2, 0, 0, LW, H, NAVY)
rect(s2, LW, 0, W - LW, H, DGREY)
rect(s2, 0, 0, LW, Inches(0.08), RED)

slide_number(s2, 2)
tag(s2, "THE SOLUTION", Inches(0.45), Inches(0.7))

txt(s2, "Refugee Assistant Switzerland",
    Inches(0.45), Inches(1.3), Inches(4.9), Inches(0.65),
    size=26, bold=True, color=WHITE)

txt(s2, "A free AI chatbot available 24/7 — answering questions about\n"
        "asylum, work, education, healthcare, and integration\n"
        "in the user's own language.",
    Inches(0.45), Inches(2.1), Inches(4.9), Inches(1.0),
    size=13, color=BGREY)

red_line(s2, Inches(3.25))

bullet_list(s2, [
    "Multilingual — 15+ languages (Arabic, Tigrinya, Dari, French…)",
    "Covers everything — work rights, school, health, permits, appeals",
    "Grounded in 27 verified Swiss official sources (SEM, OSAR, ch.ch)",
    "Always on — no appointment, no waiting, no account needed",
    "Proactively warns about critical deadlines (e.g. 30-day appeal)",
    "Free for every refugee in Switzerland",
], Inches(0.45), Inches(3.45), Inches(4.9),
   text_color=WHITE, dot_color=RED, size=13)

# Right — phone mockup
px = LW + Inches(1.5)
py = Inches(0.45)
pw = Inches(3.5)
ph = Inches(6.4)

rect(s2, px, py, pw, ph, RGBColor(0x2a, 0x2a, 0x3e))
rect(s2, px + Inches(0.14), py + Inches(0.22),
     pw - Inches(0.28), ph - Inches(0.35), DGREY)

# App header
rect(s2, px + Inches(0.14), py + Inches(0.65),
     pw - Inches(0.28), Inches(0.7), RED)
txt(s2, "Refugee Assistant Switzerland",
    px + Inches(0.22), py + Inches(0.7),
    pw - Inches(0.4), Inches(0.32),
    size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s2, "Online · Replies in seconds",
    px + Inches(0.22), py + Inches(1.0),
    pw - Inches(0.4), Inches(0.25),
    size=8, color=BGREY, align=PP_ALIGN.CENTER)

# Chat messages
chats = [
    ("user", "What is Permit F?"),
    ("bot",  "Permit F = provisionally admitted.\nYou may work with cantonal authorisation."),
    ("user", "هل يمكنني العمل؟"),
    ("bot",  "نعم، يمكنك العمل بعد الحصول على إذن."),
    ("user", "Puis-je inscrire mon enfant à l'école?"),
    ("bot",  "Oui, tous les enfants ont le droit à l'école\nen Suisse, quel que soit le permis."),
]
cy_c = py + Inches(1.5)
for role, msg in chats:
    if role == "user":
        bx2 = px + pw - Inches(2.6)
        bc2, tc2 = RED, WHITE
    else:
        bx2 = px + Inches(0.22)
        bc2, tc2 = RGBColor(0x2a, 0x2a, 0x3e), BGREY
    bh2 = Inches(0.52) if "\n" not in msg else Inches(0.72)
    rect(s2, bx2, cy_c, Inches(2.25), bh2, bc2)
    txt(s2, msg, bx2 + Inches(0.08), cy_c + Inches(0.06),
        Inches(2.1), bh2, size=8, color=tc2)
    cy_c += bh2 + Inches(0.1)


# ═══════════════════════════════════════════════════════════
# SLIDE 3 — TECH STACK
# ═══════════════════════════════════════════════════════════
s3 = add_slide()
rect(s3, 0, 0, W, H, WHITE)
rect(s3, 0, 0, W, Inches(0.08), RED)
rect(s3, 0, Inches(1.35), W, Inches(0.03), LGREY)

slide_number(s3, 3)
tag(s3, "TECH STACK", Inches(0.45), Inches(0.38))
txt(s3, "Architecture & Technology",
    Inches(2.5), Inches(0.32), Inches(8.0), Inches(0.7),
    size=26, bold=True, color=NAVY, align=PP_ALIGN.LEFT)

# 4 architecture layer boxes across the full width
layers = [
    (BLUE,   "FRONTEND",       "Streamlit",      "Python web interface\nZero JS required"),
    (GREEN,  "BACKEND",        "Python",         "Query processing\nSession management"),
    (ORANGE, "KNOWLEDGE BASE", "SQLite + RAG",   "27 Swiss official sources\nKeyword scoring retrieval"),
    (RED,    "AI ENGINE",      "LLaMA 3.3-70b",  "Via Groq API\nFast, multilingual"),
]
bw = Inches(2.9)
bh = Inches(2.4)
gap = Inches(0.24)
total_w = 4 * bw + 3 * gap
start_x = (W - total_w) / 2
by2 = Inches(1.6)

for i, (col, category, tech, desc) in enumerate(layers):
    bx2 = start_x + i * (bw + gap)
    rect(s3, bx2, by2, bw, bh, col)
    # category label
    txt(s3, category, bx2, by2 + Inches(0.18), bw, Inches(0.38),
        size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # divider
    rect(s3, bx2 + Inches(0.3), by2 + Inches(0.58), bw - Inches(0.6),
         Inches(0.03), RGBColor(0xFF, 0xFF, 0xFF))
    # tech name
    txt(s3, tech, bx2, by2 + Inches(0.72), bw, Inches(0.55),
        size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # description
    txt(s3, desc, bx2 + Inches(0.1), by2 + Inches(1.35), bw - Inches(0.2),
        Inches(0.9), size=11, color=RGBColor(0xFF, 0xFF, 0xFF),
        align=PP_ALIGN.CENTER)
    # arrow between boxes
    if i < len(layers) - 1:
        txt(s3, "→", bx2 + bw, by2 + bh / 2 - Inches(0.25),
            gap, Inches(0.5), size=20, bold=True,
            color=BGREY, align=PP_ALIGN.CENTER)

# Bottom detail row
details = [
    ("Language Detection",  "Automatic — user writes in any language,\nbot responds in the same language"),
    ("30-Day Alert",        "Bot proactively warns about critical\nSwiss appeal deadlines"),
    ("Conversation Memory", "Full chat history stored in SQLite,\ncontext preserved across turns"),
    ("Source Citation",     "Answers grounded in SEM, OSAR, ch.ch\n— not guesses, real Swiss law"),
]
dw = Inches(2.9)
dy = by2 + bh + Inches(0.35)
dh = Inches(1.55)

for i, (title, desc) in enumerate(details):
    dx = start_x + i * (dw + gap)
    rect(s3, dx, dy, dw, dh, LGREY)
    rect(s3, dx, dy, dw, Inches(0.06), RED)
    txt(s3, title, dx + Inches(0.15), dy + Inches(0.14),
        dw - Inches(0.3), Inches(0.38),
        size=11, bold=True, color=NAVY)
    txt(s3, desc, dx + Inches(0.15), dy + Inches(0.52),
        dw - Inches(0.3), Inches(0.9),
        size=10, color=MGREY)


# ═══════════════════════════════════════════════════════════
# SLIDE 4 — HOW IT WORKS (RAG) — Architecture Defense bonus
# ═══════════════════════════════════════════════════════════
s4 = add_slide()
rect(s4, 0, 0, W, H, WHITE)
rect(s4, 0, 0, W, Inches(0.08), RED)
rect(s4, 0, Inches(1.35), W, Inches(0.03), LGREY)

slide_number(s4, 4)
tag(s4, "HOW IT WORKS", Inches(0.45), Inches(0.38))
txt(s4, "RAG Pipeline — End to End",
    Inches(2.5), Inches(0.32), Inches(8.0), Inches(0.7),
    size=26, bold=True, color=NAVY)

# Flow steps (horizontal)
flow = [
    (NAVY,   "1. User Input",     "Any question\nin any language"),
    (BLUE,   "2. Keyword Search", "Match against\n27 Swiss sources"),
    (ORANGE, "3. Source Inject",  "Top 3 sources\nadded to prompt"),
    (GREEN,  "4. LLM Generate",   "LLaMA 3.3-70b\nvia Groq API"),
    (RED,    "5. Answer",         "Shown in user's\nlanguage"),
]
fw = Inches(2.1)
fh = Inches(1.8)
fgap = Inches(0.28)
total_fw = len(flow) * fw + (len(flow) - 1) * fgap
fx0 = (W - total_fw) / 2
fy = Inches(1.7)

for i, (col, title, desc) in enumerate(flow):
    fx = fx0 + i * (fw + fgap)
    rect(s4, fx, fy, fw, fh, col)
    txt(s4, title, fx, fy + Inches(0.22), fw, Inches(0.45),
        size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(s4, fx + Inches(0.25), fy + Inches(0.7),
         fw - Inches(0.5), Inches(0.03), WHITE)
    txt(s4, desc, fx, fy + Inches(0.9), fw, Inches(0.75),
        size=11, color=WHITE, align=PP_ALIGN.CENTER)
    if i < len(flow) - 1:
        txt(s4, "→", fx + fw, fy + fh / 2 - Inches(0.25),
            fgap, Inches(0.5), size=18, bold=True,
            color=BGREY, align=PP_ALIGN.CENTER)

# Example query banner
rect(s4, Inches(1.5), fy + fh + Inches(0.4),
     W - Inches(3.0), Inches(0.75), RGBColor(0xFF, 0xF5, 0xF5))
txt(s4, 'Example: User asks "Can I work with Permit F?" in Arabic  →  '
        'Bot finds SEM source  →  Responds in Arabic with correct answer  →  Saves to history',
    Inches(1.7), fy + fh + Inches(0.53),
    W - Inches(3.4), Inches(0.45),
    size=12, color=RED, italic=True, align=PP_ALIGN.CENTER)

# Design decisions section
txt(s4, "Key Architecture Decisions",
    Inches(0.5), fy + fh + Inches(1.5), W - Inches(1.0), Inches(0.45),
    size=16, bold=True, color=NAVY)

decisions = [
    ("Why keyword search, not vectors?",
     "Swiss asylum domain is structured & predictable. Keyword scoring is faster,\n"
     "cheaper, and accurate enough for this domain — no GPU needed."),
    ("Why Groq + LLaMA 3.3-70b?",
     "Fast inference, multilingual capability, OpenAI-compatible API.\n"
     "Free tier covers MVP usage — easy to swap for another model later."),
    ("Why SQLite?",
     "Zero-config, serverless, perfect for single-user Streamlit deployment.\n"
     "Stores both the knowledge base and conversation history in one file."),
]
dw2 = Inches(3.9)
dy2 = fy + fh + Inches(2.1)
dh2 = Inches(1.4)
dx2_0 = Inches(0.5)

for i, (q, a) in enumerate(decisions):
    dx2 = dx2_0 + i * (dw2 + Inches(0.25))
    rect(s4, dx2, dy2, dw2, dh2, LGREY)
    rect(s4, dx2, dy2, dw2, Inches(0.06), RED)
    txt(s4, q, dx2 + Inches(0.15), dy2 + Inches(0.14),
        dw2 - Inches(0.3), Inches(0.38),
        size=11, bold=True, color=NAVY)
    txt(s4, a, dx2 + Inches(0.15), dy2 + Inches(0.52),
        dw2 - Inches(0.3), Inches(0.82),
        size=10, color=MGREY)


# ── Save ─────────────────────────────────────────────────
out = "Refugee_Assistant_PitchDeck.pptx"
prs.save(out)
print(f"Saved: {out}")
